"""PPTX text extraction with per-slide [PAGE N] markers."""

from __future__ import annotations

from pptx import Presentation


def parse_pptx(file_path: str) -> str:
    """Return text of a .pptx file with [PAGE N] markers per slide."""
    prs = Presentation(file_path)
    slides: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        slides.append(f"[PAGE {index}]\n" + "\n".join(texts))
    return "\n\n".join(slides)
